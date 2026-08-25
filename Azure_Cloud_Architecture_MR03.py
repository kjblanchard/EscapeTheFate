#!/usr/bin/env python3
"""
Generate Azure Cloud Architecture PowerPoint - MR03 equivalent
Modeled after the AWS MR03 Cloud Architecture presentation structure
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

# Constants for Ally branding
ALLY_PURPLE = RGBColor(0x82, 0x00, 0x7E)
ALLY_DARK = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_BLUE = RGBColor(0x00, 0x3B, 0x71)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def add_title_bar(slide, title_text):
    """Add a colored title bar at the top of a slide"""
    # Title background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ALLY_PURPLE
    shape.line.fill.background()

    # Title text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)


def add_subtitle(slide, text, top=1.3):
    """Add a subtitle below the title bar"""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ALLY_DARK


def add_body_text(slide, text, left=0.5, top=2.0, width=12, height=5, font_size=14):
    """Add body text to a slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ALLY_DARK
        p.space_after = Pt(6)


def add_bullet_points(slide, bullets, left=0.5, top=2.0, width=12, height=5, font_size=13):
    """Add bulleted text"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = ALLY_DARK
        p.space_after = Pt(8)
        p.level = 0


def add_qa_slide(question, answer_bullets, section_title="Architecture Questions"):
    """Add a Q&A formatted slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_title_bar(slide, section_title)

    # Question box
    q_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.0)
    )
    q_shape.fill.solid()
    q_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE0, 0xF0)
    q_shape.line.fill.background()
    tf = q_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = f"Q: {question}"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ALLY_PURPLE

    # Answer
    add_bullet_points(slide, answer_bullets, top=2.6, font_size=12)
    return slide


def add_diagram_placeholder(slide, text, left=1.0, top=3.0, width=11, height=3.5):
    """Add a placeholder box indicating where a diagram from Confluence would go"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF)
    shape.line.color.rgb = AZURE_BLUE
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"[Diagram: {text}]"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = AZURE_BLUE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Insert from Confluence attachment"
    p2.font.size = Pt(11)
    p2.font.color.rgb = MEDIUM_GRAY
    p2.alignment = PP_ALIGN.CENTER


# =============================================================================
# SLIDE 1: Title Slide
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Full purple background
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = ALLY_PURPLE
bg_shape.line.fill.background()

# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "MR03 - Cloud Architecture"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1.0))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Azure Environment"
p2.font.size = Pt(28)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

# Date
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = f"FRB Cloud Computing ECM - Objective 2 Cloud Architecture"
p3.font.size = Pt(16)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER

p4 = tf3.add_paragraph()
p4.text = datetime.date.today().strftime("%B %Y")
p4.font.size = Pt(14)
p4.font.color.rgb = WHITE
p4.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 2: Table of Contents
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Table of Contents")
toc_items = [
    "1. Azure Environment Overview",
    "2. Networking & Connectivity",
    "3. Identity & Access Management",
    "4. Policy Governance",
    "5. Application Log Flow & Monitoring",
    "6. Hybrid DNS",
    "7. GitLab Integration & CI/CD",
    "8. NAOTest Tenant Architecture",
    "9. Architecture Questions & Answers",
]
add_bullet_points(slide, toc_items, top=1.5, font_size=16)

# =============================================================================
# SLIDE 3: Azure Environment Overview
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Azure Environment Overview")
add_subtitle(slide, "Tenant & Landing Zone Structure")
overview_bullets = [
    "• Ally uses the NAO tenant for all Azure cloud workloads",
    "• Dev/Test/Prod application environments are all within NAO in the production landing zone",
    "• Production landing zone has its own set of Azure Policy, Firewalls, Networking/ExpressRoute, and guardrails",
    "• Canary landing zone in NAO for Cloud Foundations and Palo Alto testing",
    "• NAOTest tenant enabled for specific use cases (Jeffrey Czerak / AVD team)",
    "• Cloud Adoption Framework (CAF) implemented using Terraform",
    "• Infrastructure as Code: github.com/Azure/terraform-azurerm-caf-enterprise-scale",
    "• Primary Region: US East 2 | Secondary Region: US Central",
]
add_bullet_points(slide, overview_bullets, top=2.0, font_size=13)

# =============================================================================
# SLIDE 4: Networking - High Level Architecture
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Azure Networking - High Level Architecture")
add_subtitle(slide, "Landing Zone Network Design")
net_bullets = [
    "• Multiple landing zones, each requiring unique regional CIDR and ExpressRoute connections via Equinix",
    "• Dedicated Palo Alto Firewalls in centralized Connectivity subscription per landing zone",
    "• Two Palo deployments: Trusted traffic routing and Untrusted traffic routing",
    "• Each child subscription is peered to the Connectivity subscription for routing",
    "• Subnets designated as 'public' or 'private' in a single VNet controlling route table alignment",
    "• All resources placed within VNets using Private Endpoints (internal firewalls required)",
    "• Private resources available only via Azure Private DNS (not publicly resolvable outside Ally network)",
    "• Public resources require public subnet + route table creating DMZ approach",
]
add_bullet_points(slide, net_bullets, top=2.0, font_size=12)

# =============================================================================
# SLIDE 5: Networking Diagram
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Azure Networking - Architecture Diagram")
add_subtitle(slide, "Single Landing Zone Instance")
add_diagram_placeholder(slide, "Azure Network - High Level / Landing Zones Diagram\n(Source: Confluence - Azure Networking Overview - image-2025-11-19_9-17-32.png)", top=2.0, height=4.5)

# =============================================================================
# SLIDE 6: Multi-Region Networking
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Multi-Region Networking")
add_subtitle(slide, "US East 2 (Primary) & US Central (Secondary)")
mr_bullets = [
    "• Primary: US East 2 | Secondary: US Central",
    "• Regional connectivity via peering of Connectivity Virtual Networks",
    "• Intra-subscription routing: two regional VNets peered directly",
    "• Cross-subscription routing: via Connectivity subscription (subject to Palo Alto inspection)",
]
add_bullet_points(slide, mr_bullets, top=2.0, font_size=13)
add_diagram_placeholder(slide, "Azure Region Networking - Merged.png\n(Multi-region peering topology)", top=3.8, height=3.0)

# =============================================================================
# SLIDE 7: Key Networking Concepts
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Key Networking Concepts")
concepts = [
    "• Virtual Networks (VNets) - Private network building blocks in Azure",
    "• Subnets - Segmented CIDR allocations within VNets",
    "• Network Security Groups (NSGs) - Inbound/outbound traffic rules",
    "• Route Tables - Custom routing overriding Azure defaults",
    "• Virtual Network Peering - Seamless VNet-to-VNet connectivity via Microsoft backbone",
    "• Private Endpoints - Private IP interfaces connecting to Azure PaaS via Private Link",
    "• Private DNS Resolver - Queries Azure DNS private zones from on-prem and vice versa",
    "• ExpressRoute - Private connectivity to Microsoft cloud via Equinix",
    "• Virtual Network Gateway - IP route exchange and traffic routing for ExpressRoute",
    "• Service Endpoints - Secure direct connectivity to Azure services over backbone",
]
add_bullet_points(slide, concepts, top=1.5, font_size=12)

# =============================================================================
# SLIDE 8: Identity & Access Management - Overview
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Identity & Access Management")
add_subtitle(slide, "Azure IAM Model Overview")
iam_bullets = [
    "• Azure AD (Entra ID) as Identity Provider - synced from on-premises AD",
    "• IAM implemented in Terraform with centralized state management and drift remediation",
    "• 'Human' users: AD user with role attachments (interactive use cases)",
    "• 'Machine' users: App Registrations with role attachments (CI/CD, automation)",
    "• 'Horizontal' access: Cross-subscription (e.g., BillingReadOnly)",
    "• 'Vertical' access: Per-subscription (e.g., Developer, ReadOnly)",
    "• Permission Boundaries: 'Not Actions' on roles prevent specific privileged actions",
    "• Break-glass access: Elevated access during P1/P2/P3 incidents",
    "• Strict IP whitelist enforced via Azure AD Conditional Access",
]
add_bullet_points(slide, iam_bullets, top=2.0, font_size=12)

# =============================================================================
# SLIDE 9: IAM - Human Access Flow
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "IAM - Human Access Flow")
add_subtitle(slide, "User Authentication & Authorization")
human_bullets = [
    "• Users login via Azure Portal or AZ CLI (temporary token)",
    "• Azure AD enforces strict IP whitelist (Ally-owned network only)",
    "• Users assigned to one or many Azure AD groups",
    "• Groups have roles attached granting specific actions at specific scopes",
    "• Users synced from on-premises AD to Azure AD",
    "",
    "Role Attachment Scopes:",
    "• Management Group → Subscription → Resource Group → Resource",
    "• Scope determines Horizontal vs Vertical access",
]
add_bullet_points(slide, human_bullets, top=2.0, font_size=13)
add_diagram_placeholder(slide, "IAM Role Attachment Model\n(Source: Azure Identity and Access Management - image-2024-2-28_9-26-13.png)", top=5.0, height=2.0)

# =============================================================================
# SLIDE 10: IAM - Machine Access & Federated IAM
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "IAM - Machine Access & Federated IAM")
add_subtitle(slide, "App Registrations & Self-Service Roles")
machine_bullets = [
    "App Registration (Machine Users):",
    "• Managed in Terraform (azure-tf-app-registration)",
    "• Requires documented use case with specific permissions/scope",
    "• Keys expire after 1 year, must be rotated",
    "",
    "Federated IAM (Self-Service):",
    "• Similar to AWS Federated IAM - tailored for Azure",
    "• Self-service by app teams via azure-tf-federated-iam repo",
    "• Module validates no restricted/privileged actions added to roles",
    "• Replaces deprecated Azure AD Custom Role and Custom Group processes",
    "• Horizontal roles & highly privileged roles maintained by Cloud Foundations",
]
add_bullet_points(slide, machine_bullets, top=2.0, font_size=12)

# =============================================================================
# SLIDE 11: Policy Governance
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Azure Policy Governance")
add_subtitle(slide, "Enforcement of Configuration & Security Requirements")
policy_bullets = [
    "• Azure Policy enforces organizational standards and assesses compliance at scale",
    "• Aggregated compliance dashboard with per-resource, per-policy granularity",
    "• Bulk remediation for existing resources; automatic remediation for new resources",
    "• Managed as extension of Cloud Adoption Framework using Terraform",
    "• Source: gitlab - azure-tf-infrastructure/policy",
    "",
    "Three Core Components:",
    "• Policy Definitions - Conditions + effects (e.g., mandatory tagging)",
    "• Policy Sets (Initiatives) - Collections of policies for a singular goal",
    "• Policy Assignments - Definitions/initiatives assigned to specific scopes (Management Group → Resource)",
    "• Assignments inherited by all child resources (with subscope exclusion capability)",
]
add_bullet_points(slide, policy_bullets, top=2.0, font_size=12)

# =============================================================================
# SLIDE 12: Policy Governance Diagram
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Azure Policy - Compliance Dashboard")
add_subtitle(slide, "Aggregated View of Environment State")
add_diagram_placeholder(slide, "Azure Policy Compliance Dashboard\n(Source: Azure Policy Governance - image-2024-3-4_10-18-29.png)", top=2.0, height=4.5)

# =============================================================================
# SLIDE 13: Application Log Flow
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Application Log Flow & Monitoring")
add_subtitle(slide, "Azure Monitor → Event Hub → Splunk/Dynatrace")
log_bullets = [
    "• Azure Monitor collects log data from all Azure resources and services",
    "• Resources configured with Diagnostic Settings to send logs to Event Hub namespace",
    "• Event Hub namespace acts as centralized ingestion point with multiple Event Hubs",
    "• Each Event Hub dedicated to a specific log data stream",
    "• Data forwarded to Splunk HTTP Event Collector (HEC) endpoints",
    "• Event Hubs centralized in Landingzone-management-prod subscription",
    "• Logs backed up to S3 if Splunk unavailable",
    "",
    "Flow: Resource → Diagnostic Setting → Azure Monitor → Event Hub → Splunk HEC",
]
add_bullet_points(slide, log_bullets, top=2.0, font_size=12)

# =============================================================================
# SLIDE 14: Log Flow Diagram
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Application Log Flow - Architecture Diagram")
add_diagram_placeholder(slide, "Application Log Flow Architecture\n(Source: Azure Application Log Flow - image-2024-7-12_13-33-34.png)", top=1.5, height=5.5)

# =============================================================================
# SLIDE 15: NSG Flow Logging & Observability
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "NSG Flow Logging & Observability Transition")
add_subtitle(slide, "Network Traffic Capture & Dynatrace Migration")
nsg_bullets = [
    "NSG Flow Logging:",
    "• Network traffic captured through NSG flow logs",
    "• Flow logs sent to centralized Azure Storage account",
    "• BUs cannot create new NSGs - must use pre-configured NSGs with flow logs",
    "• Retention: 90 days",
    "• DNS log capture currently in preview",
    "",
    "Splunk (Current SIEM):",
    "• Log Analytics Workspaces send data via Event Hub automatically (Azure Policy)",
    "• Splunk retention: 90 days",
    "",
    "Dynatrace Transition:",
    "• All observability-related events transitioning to Dynatrace",
    "• Security events remain in Splunk (compliance, audit, threat detection)",
    "• Dynatrace: metrics, traces, logs for performance/availability",
]
add_bullet_points(slide, nsg_bullets, top=2.0, font_size=11)

# =============================================================================
# SLIDE 16: Hybrid DNS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Azure Hybrid DNS")
add_subtitle(slide, "Private DNS Resolver Architecture")
dns_bullets = [
    "• Azure Private DNS Resolvers - highly available, scalable DNS service",
    "• Resolves DNS requests within Azure for on-prem, Azure native zones, and external domains",
    "",
    "DNS Access Patterns:",
    "• System DNS Rules: External DNS domain resolution",
    "• Inbound Resolver Rules: VPC endpoints, cross-account private hosted zones, private links",
    "• Outbound Resolver Rules: On-prem DNS domain resolution",
    "• On-prem Azure Domains: On-prem accessing Azure private hosted zones and private link zones",
    "",
    "Design Principles:",
    "• Single pane of glass | Centralized deployment | Scalable",
    "• Source: azure-tf-infrastructure/modules/modules/dns/dns_endpoint.tf",
]
add_bullet_points(slide, dns_bullets, top=2.0, font_size=12)

# =============================================================================
# SLIDE 17: DNS Diagram
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Azure Hybrid DNS - Architecture")
add_diagram_placeholder(slide, "Azure Top Level DNS Integration\n(Source: Azure Hybrid DNS - Azure Top Level DNS Integration.jpeg)", top=1.5, height=5.5)

# =============================================================================
# SLIDE 18: GitLab Integration
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "GitLab Integration & CI/CD")
add_subtitle(slide, "Deployment Pipeline from GitLab.com to Azure")
gl_bullets = [
    "• Reuses existing GitLab.com runners hosted in AWS",
    "• Keeps runner management minimal until Azure growth warrants dedicated runners",
    "",
    "Pipeline Flow:",
    "1. GitLab pipeline targets Azure using existing runner in aligned AWS account",
    "2. PCE pipeline utility detects Azure deployment via specific parameters",
    "3. Pipeline retrieves AWS Secrets Manager value representing Azure subscription credentials",
    "4. Azure CLI invoked to assume the App Registration for that Azure subscription",
    "5. Pipeline does lookup for pre-provisioned Storage Account (Terraform state)",
    "6. Pipeline completes via AWS runner, creating infrastructure in target Azure subscription",
]
add_bullet_points(slide, gl_bullets, top=2.0, font_size=12)

# =============================================================================
# SLIDE 19: GitLab Integration Diagram
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "GitLab Integration - Architecture Diagram")
add_diagram_placeholder(slide, "Azure GitLab Integration Flow\n(Source: Azure Gitlab Integration - Azure Subscriptions and RBAC(2).png)", top=1.5, height=5.5)

# =============================================================================
# SLIDE 20: NAOTest Tenant
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "NAOTest Tenant Architecture")
add_subtitle(slide, "Non-Production Tenant for Testing & AVD")
nao_bullets = [
    "• NAOTest tenant enabled for Cloud Foundations testing and AVD team iteration",
    "• Allows different tenant-level settings without affecting NAO production",
    "• Use cases: Conditional Access policies (Jeffrey Czerak), Group Policy/AD testing (AVD team)",
    "• Requires standardized guardrails, subscription automation, and private connectivity",
    "",
    "Traffic Inspection:",
    "• Outbound: Inspection through Azure Firewalls",
    "• Cross-subscription: Inspection through Azure Firewalls",
    "• On-prem: Inspection through On-Ramp Firewalls",
    "",
    "Selected Approach: Option 2 (Long-term) - Separate ExpressRoute, decommission Canary",
    "• Completely separate environment from security perspective",
    "• Standardized landing zones | Consolidates Canary into NAOTest",
]
add_bullet_points(slide, nao_bullets, top=2.0, font_size=12)

# =============================================================================
# SLIDE 21: NAOTest Options Diagram
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "NAOTest - Selected Architecture (Option 2 Long-term)")
add_subtitle(slide, "Canary Decommissioned, NAOTest with Separate ExpressRoute")
add_diagram_placeholder(slide, "NAOTest Option 2 Long-term Diagram\n(Source: Azure - NAOTest - image-2025-11-13_9-46-0.png)", top=2.0, height=4.5)

# =============================================================================
# ARCHITECTURE QUESTIONS SECTION
# =============================================================================

# Q1: Network Architecture
add_qa_slide(
    "How is the network architecture designed in Azure?",
    [
        "• Hub-and-spoke model using Cloud Adoption Framework (CAF) with Terraform",
        "• Centralized Connectivity subscription with dedicated Palo Alto Firewalls",
        "• Two firewall deployments: Trusted (internal) and Untrusted (DMZ/public) traffic",
        "• Child subscriptions peered to Connectivity subscription for all routing",
        "• Subnets classified as 'public' or 'private' controlling route table alignment",
        "• All PaaS resources accessed via Private Endpoints only",
        "• Multi-region: US East 2 (Primary) peered to US Central (Secondary)",
        "• Updated architecture maintained in LucidChart",
    ],
    "Architecture Questions - Networking"
)

# Q2: Traffic Segmentation
add_qa_slide(
    "How is traffic segmented and inspected?",
    [
        "• Public subnets route through Untrusted Palo Alto firewalls (DMZ approach)",
        "• Private subnets route through Trusted Palo Alto firewalls",
        "• Cross-subscription traffic inspected via Connectivity subscription Palo firewalls",
        "• On-prem traffic inspected through On-Ramp firewalls",
        "• NSGs control inbound/outbound traffic per subnet (managed, not user-creatable)",
        "• Route tables enforce traffic direction to appropriate firewall deployment",
        "• ExpressRoute provides private connectivity (no public internet traversal)",
    ],
    "Architecture Questions - Networking"
)

# Q3: IAM
add_qa_slide(
    "How is identity and access management implemented?",
    [
        "• Azure AD (Entra ID) as sole Identity Provider, synced from on-prem AD",
        "• Strict IP whitelist via Conditional Access (Ally network only)",
        "• Human access: AD users → Azure AD Groups → Role Attachments at defined scopes",
        "• Machine access: App Registrations with scoped role attachments",
        "• Federated IAM (self-service): azure-tf-federated-iam repo with guardrails",
        "• Module validation prevents adding restricted/privileged actions to custom roles",
        "• Break-glass access available for P1/P2/P3 incidents",
        "• All IAM managed as Infrastructure as Code (Terraform) with drift remediation",
    ],
    "Architecture Questions - IAM"
)

# Q4: Horizontal vs Vertical
add_qa_slide(
    "What is the difference between Horizontal and Vertical access?",
    [
        "• Horizontal: Access across ALL subscriptions in the Azure tenant",
        "  - Example: BillingReadOnly for billing team (read all subscription costs)",
        "  - Managed in Cloud Foundations Infrastructure repo",
        "",
        "• Vertical: Access within a SPECIFIC subscription (per-account basis)",
        "  - Standard roles created during subscription creation (ReadOnly, Developer)",
        "  - Additional vertical roles created via groups pipeline (e.g., OpenAI Contributor)",
        "  - Managed through approvals with minimal permissions best practice",
        "",
        "• Users/machines can have Horizontal, Vertical, or combination of both",
        "• Role Attachment scope determines access type: Management Group (horizontal) vs Subscription/RG (vertical)",
    ],
    "Architecture Questions - IAM"
)

# Q5: Policy Governance
add_qa_slide(
    "How are Azure policies managed and enforced?",
    [
        "• Azure Policy service enforces configuration and security requirements in Landing Zones",
        "• Managed via Cloud Adoption Framework extension using Terraform",
        "• Source: azure-tf-infrastructure/policy (GitLab)",
        "",
        "• Policy Definitions: Conditions + effects (e.g., mandatory resource tagging)",
        "• Policy Sets (Initiatives): Collections for overarching goals (e.g., Defender for Cloud monitoring)",
        "• Policy Assignments: Applied at Management Group level, inherited by all children",
        "",
        "• Compliance dashboard provides aggregated per-resource, per-policy view",
        "• Bulk remediation for existing resources; automatic remediation for new resources",
        "• Some resources auto-log to Splunk via policy (see azadvertizer.net for supported list)",
    ],
    "Architecture Questions - Policy"
)

# Q6: Logging
add_qa_slide(
    "How are application logs collected and forwarded?",
    [
        "• Azure Monitor collects diagnostic logs from all Azure resources automatically",
        "• Diagnostic Settings configured to forward logs to centralized Event Hub namespace",
        "• Event Hub namespace in Landingzone-management-prod subscription",
        "• Multiple Event Hubs per namespace (one per log data stream type)",
        "• Event Hubs format data and forward to Splunk HEC endpoints and Dynatrace",
        "• Logs backed up to S3 if Splunk is unavailable",
        "",
        "• NSG Flow Logs: Captured and sent to centralized Azure Storage (90-day retention)",
        "• Splunk remains primary SIEM (security events, 90-day retention)",
        "• Dynatrace: Transitioning all observability events (metrics, traces, perf logs)",
        "• DNS logging: Currently in preview, being evaluated with Azure support",
    ],
    "Architecture Questions - Logging"
)

# Q7: DNS
add_qa_slide(
    "How does hybrid DNS work between on-premises and Azure?",
    [
        "• Azure Private DNS Resolvers provide highly available, scalable DNS resolution",
        "• Centralized deployment model with single pane of glass management",
        "",
        "DNS Resolution Patterns:",
        "• System DNS Rules → External domain resolution",
        "• Inbound Resolver → VNet endpoints, cross-account private hosted zones, private links",
        "• Outbound Resolver → On-prem DNS domain resolution",
        "• On-prem → Azure: Accesses Azure private hosted zones and private link zones",
        "",
        "• All private resources resolvable only via Azure Private DNS",
        "• Not publicly resolvable outside Ally network",
        "• Managed in Terraform: dns_endpoint.tf module",
    ],
    "Architecture Questions - DNS"
)

# Q8: CI/CD
add_qa_slide(
    "How do CI/CD pipelines deploy to Azure?",
    [
        "• GitLab.com pipelines leverage existing AWS-hosted runners",
        "• PCE pipeline utility detects Azure target via specific parameters",
        "• Authentication flow:",
        "  1. Pipeline retrieves Azure credentials from AWS Secrets Manager",
        "  2. Azure CLI assumes App Registration in target subscription",
        "  3. Pre-provisioned Storage Account discovered for Terraform state",
        "  4. Infrastructure created in target Azure subscription",
        "",
        "• Each subscription has a pre-provisioned Storage Account for Terraform state",
        "• App Registrations per subscription with specific scoped permissions",
        "• Approach reuses existing infrastructure until Azure growth warrants dedicated runners",
    ],
    "Architecture Questions - CI/CD"
)

# Q9: NAOTest
add_qa_slide(
    "What is the NAOTest tenant and why does it exist?",
    [
        "• NAOTest is a separate Azure tenant for non-production testing",
        "• Enables different tenant-level settings without affecting NAO production",
        "",
        "Use Cases:",
        "• Cloud Foundations: Test Azure policy, routing, firewall changes before prod",
        "• AVD Team: Faster iteration on group policy, conditional access policies",
        "• Jeffrey Czerak: Different conditional access policies than production",
        "",
        "Selected Architecture (Option 2 Long-term - 11/13/25 decision):",
        "• Separate ExpressRoute connection for complete security isolation",
        "• Decommission Canary landing zone and consolidate into NAOTest",
        "• Standardized landing zone approach (like-for-like with production)",
        "• Reduces costs by eliminating redundant Canary environment",
    ],
    "Architecture Questions - NAOTest"
)

# Q10: Security boundaries
add_qa_slide(
    "What security boundaries and guardrails are in place?",
    [
        "• Azure Policy: Enforces configuration standards at Management Group level",
        "• Palo Alto Firewalls: Dedicated trusted/untrusted traffic inspection",
        "• NSGs: Network-level traffic control (managed by Cloud Foundations, not app teams)",
        "• Private Endpoints: All PaaS resources accessed privately only",
        "• Azure AD Conditional Access: IP whitelist restricting to Ally network",
        "• Role 'Not Actions': Prevent privileged operations even with wildcard permissions",
        "• Federated IAM validation: Blocks restricted actions in self-service role creation",
        "• App Registration key rotation: 1-year expiration enforced",
        "• NSG Flow Logs: 90-day retention for network traffic audit",
        "• Splunk SIEM: Security event ingestion and compliance monitoring",
    ],
    "Architecture Questions - Security"
)

# =============================================================================
# SLIDE: References
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "References & Source Documentation")
refs = [
    "Confluence Pages:",
    "• Azure Policy Governance - https://ally-financial.atlassian.net/wiki/spaces/CLOUD/pages/404754781",
    "• Azure Networking Overview - https://ally-financial.atlassian.net/wiki/spaces/CLOUD/pages/404754683",
    "• Azure Application Log Flow - https://ally-financial.atlassian.net/wiki/spaces/CLOUD/pages/404751131",
    "• Azure Hybrid DNS - https://ally-financial.atlassian.net/wiki/spaces/CLOUD/pages/404755127",
    "• Azure GitLab Integration - https://ally-financial.atlassian.net/wiki/spaces/CLOUD/pages/404754691",
    "• Azure Identity and Access Management - https://ally-financial.atlassian.net/wiki/spaces/CLOUD/pages/404751077",
    "• Azure NAOTest - https://ally-financial.atlassian.net/wiki/spaces/CLOUD/pages/404758323",
    "",
    "GitLab Repositories:",
    "• azure-tf-infrastructure (Policy, DNS, Networking)",
    "• azure-tf-app-registration (Machine Users)",
    "• azure-tf-federated-iam (Self-Service Roles)",
    "• azure-tf-groups (Deprecated - Azure AD Groups)",
]
add_bullet_points(slide, refs, top=1.5, font_size=11)

# Save
output_path = "/Users/qzlcnf/git/external/EscapeTheFate/Azure_MR03_Cloud_Architecture.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
